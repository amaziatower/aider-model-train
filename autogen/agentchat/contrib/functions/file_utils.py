from .functions_utils import requires_python_packages, requires_secret


@requires_python_packages("pdfminer.six", "requests")
def read_text_from_pdf(file_path: str) -> str:
    """
    Reads text from a PDF file and returns it as a string.

    Args:
        file_path (str): The path to the PDF file.

    Returns:
        str: The extracted text from the PDF file.
    """
    import io
    import requests
    from pdfminer.high_level import PDFResourceManager, PDFPageInterpreter
    from pdfminer.converter import TextConverter
    from pdfminer.pdfpage import PDFPage

    resource_manager = PDFResourceManager()
    text_stream = io.StringIO()
    converter = TextConverter(resource_manager, text_stream)
    interpreter = PDFPageInterpreter(resource_manager, converter)

    if file_path.startswith("http://") or file_path.startswith("https://"):
        response = requests.get(file_path)
        file = io.BytesIO(response.content)
    else:
        file = open(file_path, "rb")

    for page in PDFPage.get_pages(file):
        interpreter.process_page(page)

    text = text_stream.getvalue()
    converter.close()
    text_stream.close()

    return text


@requires_python_packages("python-docx")
def read_text_from_docx(file_path: str) -> str:
    """
    Reads text from a DOCX file and returns it as a string.

    Args:
        file_path (str): The path to the DOCX file.

    Returns:
        str: The extracted text from the DOCX file.
    """
    from docx import Document

    doc = Document(file_path)
    paragraphs = [p.text for p in doc.paragraphs]
    text = "\n".join(paragraphs)

    return text


@requires_python_packages("pillow", "requests", "easyocr")
def read_text_from_image(file_path: str) -> str:
    """
    Reads text from an image file or URL and returns it as a string.

    Args:
        file_path (str): The path to the image file or URL.

    Returns:
        str: The extracted text from the image file or URL.
    """
    import io
    import requests
    import easyocr
    from PIL import Image

    reader = easyocr.Reader(["en"])  # specify the language(s)

    if file_path.startswith("http://") or file_path.startswith("https://"):
        response = requests.get(file_path)
        image = Image.open(io.BytesIO(response.content))
    else:
        image = Image.open(file_path)

    output = reader.readtext(image)

    # The output is a list of tuples, each containing the coordinates of the text and the text itself.
    # We join all the text pieces together to get the final text.
    text = " ".join([item[1] for item in output])

    return text


@requires_python_packages("python-pptx")
def read_text_from_pptx(file_path: str) -> str:
    """
    Reads text from a PowerPoint file and returns it as a string.

    Args:
        file_path (str): The path to the PowerPoint file.

    Returns:
        str: The extracted text from the PowerPoint file.
    """
    from pptx import Presentation

    presentation = Presentation(file_path)
    text = ""

    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                text += shape.text

    return text


@requires_python_packages("pandas")
def read_text_from_xlsx(file_path: str) -> str:
    """
    Reads text from an Excel file and returns it as a string.

    Args:
        file_path (str): The path to the Excel file.

    Returns:
        str: The extracted text from the Excel file.
    """
    import pandas as pd

    df = pd.read_excel(file_path)
    text = df.to_string(index=False)

    return text


@requires_python_packages("SpeechRecognition", "requests")
def read_text_from_audio(file_path: str) -> str:
    """
    Reads text from an audio file or a URL and returns it as a string.

    Args:
        file_path (str): The path to the audio file or the URL.

    Returns:
        str: The extracted text from the audio file or the URL.
    """
    import requests
    import speech_recognition as sr
    import tempfile

    recognizer = sr.Recognizer()

    if file_path.startswith("http"):
        response = requests.get(file_path)
        with tempfile.NamedTemporaryFile(delete=True, suffix=".wav") as temp_audio:
            temp_audio.write(response.content)
            temp_audio.seek(0)
            with sr.AudioFile(temp_audio.name) as source:
                audio = recognizer.record(source)
    else:
        with sr.AudioFile(file_path) as source:
            audio = recognizer.record(source)

    text = recognizer.recognize_google(audio)

    return text


@requires_secret("OPENAI_API_KEY")
@requires_python_packages("openai")
def caption_image_using_gpt4v(file_path_or_url: str, prompt: str = "What's in this image?") -> str:
    """
    Generates a caption for an image using the GPT-4 Vision model from OpenAI.

    Args:
        file_path_or_url (str): The path to the image file or the URL.
        prompt (str, Optional): The prompt to pass to GPT-v

    Returns:
        str: The caption generated for the image.
    """
    import os
    import openai
    from openai import OpenAI

    caption = ""

    openai.api_key = os.environ["OPENAI_API_KEY"]
    client = OpenAI()

    # check if the file_path_or_url is a URL
    if file_path_or_url.startswith("http://") or file_path_or_url.startswith("https://"):
        image_url = file_path_or_url

        response = client.chat.completions.create(
            model="gpt-4-vision-preview",
            messages=[
                {
                    "role": "user",
                    "content": [
                        {"type": "text", "text": prompt},
                        {
                            "type": "image_url",
                            "image_url": {
                                "url": image_url,
                            },
                        },
                    ],
                }
            ],
            max_tokens=300,
        )
        caption = response.choices[0]
    else:
        caption = "Please provide a valid image URL"
    return caption
